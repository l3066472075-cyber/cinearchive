"""导出 PPT：把影领家的 5 问 + 推荐影片，生成一份观电影法风格的带领方案 PPT。"""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# 品牌色（禅说电影金棕调）
BG = RGBColor(0x10, 0x0E, 0x0C)
PANEL = RGBColor(0x1A, 0x16, 0x12)
GOLD = RGBColor(0xC9, 0xA1, 0x5C)
GOLD_SOFT = RGBColor(0xE5, 0xC9, 0x8F)
INK = RGBColor(0xEC, 0xE3, 0xD3)
INK_2 = RGBColor(0xB9, 0xAD, 0x99)

FONT = "Microsoft YaHei"

QUOTE = "生命是条长河，最终渡你的还是自己"


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _box(slide, left, top, width, height, fill=PANEL):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _text(slide, left, top, width, height, lines, size=16, color=INK, bold=False, align=None):
    from pptx.enum.text import PP_ALIGN

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = FONT
        if align:
            p.alignment = align
    return tb


def _slide_title(slide, kicker, title):
    _text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4), kicker, size=13, color=GOLD)
    _text(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(0.7), title, size=28, color=INK, bold=True)


def build_ppt(answers: dict, movies: list[dict]) -> bytes:
    """生成 PPTX 字节。movies 每项含 title/synopsis/therapy_notes/trigger_warnings/discussion_questions。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = answers.get("theme") or answers.get("value") or "观电影法 · 观影共修"
    audience = answers.get("audience") or ""

    # 1) 封面
    s = _blank(prs)
    _text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5), "观电影法 · 带领方案", size=16, color=GOLD)
    _text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.6), theme, size=44, color=INK, bold=True)
    _text(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.6), f"服务对象：{audience}" if audience else "", size=16, color=INK_2)
    _text(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6), f"「{QUOTE}」", size=15, color=GOLD_SOFT)

    # 2) 观影概览（5 问）
    s = _blank(prs)
    _slide_title(s, "OVERVIEW", "本次观影 · 五问概览")
    labels = {"emotion": "需求", "situation": "目标", "value": "想法", "audience": "对象", "theme": "主题"}
    lines = [f"· {labels.get(k, k)}：{v}" for k, v in answers.items() if v]
    _box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.4))
    _text(s, Inches(1.1), Inches(2.2), Inches(11.1), Inches(3.8), lines or ["（未填写）"], size=20, color=INK)

    # 3) 选片建议
    s = _blank(prs)
    _slide_title(s, "FILMS", "选片建议")
    _text(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.5), f"首选：《{movies[0]['title']}》" if movies else "暂无", size=22, color=GOLD_SOFT, bold=True)
    if len(movies) > 1:
        _text(s, Inches(0.7), Inches(2.6), Inches(11.9), Inches(0.5), "备选：" + "、".join(f"《{m['title']}》" for m in movies[1:4]), size=16, color=INK_2)

    # 4-6) 每部影片
    for i, m in enumerate(movies[:3]):
        s = _blank(prs)
        _slide_title(s, f"FILM {i + 1}", f"《{m.get('title', '')}》")
        _box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(2.6))
        _text(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(2.2),
              [m.get("synopsis", "") or "", f"治疗要点：{m.get('therapy_notes', '') or ''}"], size=15, color=INK)
        qs = m.get("discussion_questions") or []
        if qs:
            _text(s, Inches(0.7), Inches(4.8), Inches(11.9), Inches(0.4), "讨论问题", size=15, color=GOLD, bold=True)
            _text(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.6),
                  [f"{j + 1}. {q}" for j, q in enumerate(qs[:4])], size=14, color=INK_2)

    # 7) 带领流程 SOP
    s = _blank(prs)
    _slide_title(s, "SOP", "带领流程")
    steps = [
        ("破冰", "5 min", "一句自我介绍的引子，让成员放松、彼此认识"),
        ("观影", "60 min", "安静观影，聚焦「看到什么、感受到什么」"),
        ("讨论引导", "20 min", "抛讨论问题，分组分享「哪个瞬间最触动你」"),
        ("收尾", "5 min", "每人一句话收束，把感受带回生活"),
    ]
    top = 1.9
    for name, dur, desc in steps:
        _box(s, Inches(0.7), Inches(top), Inches(2.2), Inches(1.0))
        _text(s, Inches(0.9), Inches(top + 0.15), Inches(1.8), Inches(0.6), [name, dur], size=15, color=GOLD, bold=True)
        _text(s, Inches(3.1), Inches(top + 0.15), Inches(9.5), Inches(0.7), desc, size=15, color=INK)
        top += 1.15

    # 8) 触发风险
    s = _blank(prs)
    _slide_title(s, "CARE", "触发风险与伦理提醒")
    warns = []
    for m in movies[:3]:
        for w in m.get("trigger_warnings") or []:
            warns.append(f"· 《{m['title']}》：{w}")
    _box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.6))
    _text(s, Inches(1.1), Inches(2.2), Inches(11.1), Inches(3.0),
          warns or ["· 本次影片整体温和，仍请关注成员情绪反应。"], size=15, color=INK)
    _text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.8),
          "带领分寸：不越界、不诊断、不替成员下结论，只陪伴与照见。", size=14, color=GOLD)

    # 9) 结束语
    s = _blank(prs)
    _text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.0), "观电影 · 观心 · 观自己", size=26, color=GOLD_SOFT, bold=True)
    _text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.0), f"「{QUOTE}」", size=20, color=INK)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
